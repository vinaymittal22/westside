"""
Burnt Toast × TATA — AI Fashion Stylist
Business Pitch Deck  |  15 Slides  |  16:9 Widescreen
Audience: Non-technical, business-oriented panel
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Brand Palette ─────────────────────────────────────────────────
CREAM       = RGBColor(0xF0, 0xEB, 0xE0)
CREAM_SOFT  = RGBColor(0xF8, 0xF4, 0xEE)
INK         = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT      = RGBColor(0xB8, 0x49, 0x2C)   # burnt orange
ACCENT_LITE = RGBColor(0xD4, 0x6A, 0x48)
SAGE        = RGBColor(0x74, 0x8B, 0x6A)
GOLD        = RGBColor(0xC9, 0x96, 0x2E)
MUTED       = RGBColor(0x8A, 0x87, 0x82)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x12, 0x12, 0x12)
LIGHT_GREY  = RGBColor(0xE8, 0xE4, 0xDC)

W = Inches(13.33)   # slide width  (widescreen 16:9)
H = Inches(7.5)     # slide height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # fully blank layout

# ── Low-level helpers ────────────────────────────────────────────

def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, radius=0.1):
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_multiline(slide, lines, x, y, w, h,
                  font_size=18, bold=False, color=INK,
                  align=PP_ALIGN.LEFT, line_spacing=1.2, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox

def add_divider(slide, x, y, w, color=ACCENT, thickness=0.03):
    shape = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def card(slide, x, y, w, h, bg=WHITE):
    shape = slide.shapes.add_shape(5,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = LIGHT_GREY
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.05
    return shape

def dot_label(slide, icon, label, sublabel, x, y, icon_size=28, label_size=13, sub_size=11):
    """Icon + label + sublabel stacked"""
    add_text(slide, icon, x, y, 1.2, 0.5, font_size=icon_size, align=PP_ALIGN.CENTER, color=ACCENT)
    add_text(slide, label, x, y+0.55, 1.2, 0.35, font_size=label_size, bold=True,
             align=PP_ALIGN.CENTER, color=INK)
    if sublabel:
        add_text(slide, sublabel, x-0.1, y+0.9, 1.4, 0.4, font_size=sub_size,
                 align=PP_ALIGN.CENTER, color=MUTED)

def slide_header(slide, title, subtitle=None):
    """Reusable top bar for interior slides"""
    add_rect(slide, 0, 0, 13.33, 1.1, INK)
    add_text(slide, title, 0.5, 0.18, 10, 0.55,
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.72, 9, 0.35,
                 font_size=13, color=ACCENT_LITE, align=PP_ALIGN.LEFT, italic=True)
    # small brand tag top-right
    add_text(slide, "BURNT TOAST  ×  TATA", 10.2, 0.3, 3, 0.35,
             font_size=9, color=MUTED, align=PP_ALIGN.RIGHT, bold=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)                       # full black bg
add_rect(sl, 0, 0, 4.5, 7.5, ACCENT)                      # left orange panel

# Left panel text
add_text(sl, "B", 0.35, 0.5, 4, 1.2, font_size=90, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "BURNT\nTOAST", 0.4, 1.4, 3.8, 1.6, font_size=38, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "× TATA", 0.4, 3.0, 3.8, 0.5, font_size=18, color=WHITE,
         align=PP_ALIGN.LEFT, italic=True)
add_divider(sl, 0.4, 3.7, 3.4, WHITE, 0.04)
add_text(sl, "AI FASHION\nSTYLIST", 0.4, 3.9, 3.8, 1.4,
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Testing & Demo Phase  |  2026", 0.4, 6.8, 3.8, 0.45,
         font_size=11, color=RGBColor(0xFF,0xCC,0xAA), align=PP_ALIGN.LEFT)

# Right panel
add_text(sl, "Dress Smarter.\nShop Better.\nFeel Confident.",
         5.0, 1.5, 7.8, 2.5, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_divider(sl, 5.0, 4.2, 6.5, ACCENT, 0.04)
add_text(sl, "An AI-powered personal stylist built into Burnt Toast's\nshopping experience — designed for Gen Z, built for growth.",
         5.0, 4.4, 7.8, 1.2, font_size=16, color=MUTED, align=PP_ALIGN.LEFT)
add_text(sl, "Meet  TOASTIE  →", 5.0, 5.9, 5, 0.6,
         font_size=20, bold=True, color=ACCENT_LITE, align=PP_ALIGN.LEFT)

# bottom strip
add_rect(sl, 0, 7.1, 13.33, 0.4, RGBColor(0x25, 0x25, 0x25))
add_text(sl, "CONFIDENTIAL  |  BUSINESS PRESENTATION  |  TATA INTERNAL", 0, 7.15, 13.33, 0.3,
         font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA / FLOW
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "What We'll Cover Today")

items = [
    ("01", "The Problem",        "Why Gen Z needs something different"),
    ("02", "The Solution",       "What Toastie does & how it works"),
    ("03", "Key Features",       "Core capabilities of the platform"),
    ("04", "Business Impact",    "Engagement, conversions & revenue"),
    ("05", "Why It's Different", "Toastie vs. a regular chatbot"),
    ("06", "Future Roadmap",     "Scalability & what's coming next"),
    ("07", "Market Opportunity", "Gen Z + fashion e-commerce growth"),
    ("08", "Next Steps",         "From demo to full product launch"),
]

cols = [(0.4, 6.2), (6.8, 6.2)]
for i, (num, title, sub) in enumerate(items):
    col_i = i % 2
    row_i = i // 2
    x = cols[col_i][0]
    y = 1.4 + row_i * 1.35
    card(sl, x, y, cols[col_i][1], 1.1, WHITE)
    add_rect(sl, x, y, 0.55, 1.1, ACCENT)
    add_text(sl, num, x, y+0.3, 0.55, 0.5, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.65, y+0.1, 5.2, 0.45, font_size=16, bold=True, color=INK)
    add_text(sl, sub, x+0.65, y+0.55, 5.2, 0.45, font_size=12, color=MUTED)

add_text(sl, "A focused, business-first walkthrough — no technical jargon, just clear value.",
         0.4, 7.1, 12.5, 0.35, font_size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)
add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
add_text(sl, "The Problem", 0.5, 0.18, 8, 0.55, font_size=26, bold=True, color=WHITE)
add_text(sl, "Why today's online shopping experience is broken for Gen Z",
         0.5, 0.72, 9, 0.35, font_size=13, color=WHITE, italic=True)
add_text(sl, "BURNT TOAST  ×  TATA", 10.2, 0.3, 3, 0.35,
         font_size=9, color=RGBColor(0xCC,0x88,0x60), align=PP_ALIGN.RIGHT, bold=True)

problems = [
    ("😵", "Decision Fatigue",
     "Hundreds of products, zero guidance.\nShoppers don't know what works together."),
    ("🛒", "High Cart Abandonment",
     "Over 70% of carts are abandoned.\nMost shoppers leave without buying anything."),
    ("🤖", "Generic Recommendations",
     "\"You may also like\" doesn't cut it.\nGen Z wants style, not a list of products."),
    ("📱", "No Personal Stylist",
     "In-store styling = premium experience.\nOnline shopping has never matched that — until now."),
]

for i, (icon, title, body) in enumerate(problems):
    x = 0.4 + i * 3.2
    card(sl, x, 1.4, 3.0, 4.8, RGBColor(0x28, 0x28, 0x28))
    add_rect(sl, x, 1.4, 3.0, 0.08, ACCENT)
    add_text(sl, icon, x, 1.55, 3.0, 0.8, font_size=32, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(sl, title, x+0.1, 2.45, 2.8, 0.5, font_size=15, bold=True,
             color=ACCENT_LITE, align=PP_ALIGN.CENTER)
    add_divider(sl, x+0.4, 3.05, 2.2, MUTED, 0.02)
    add_text(sl, body, x+0.1, 3.2, 2.8, 1.8, font_size=12.5, color=RGBColor(0xCC,0xCC,0xCC),
             align=PP_ALIGN.CENTER)

add_rect(sl, 0, 6.4, 13.33, 0.9, RGBColor(0xB8,0x49,0x2C))
add_text(sl, "The result: lost sales, low loyalty, and a frustrated Gen Z customer who shops elsewhere.",
         0.5, 6.5, 12.3, 0.6, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — THE SOLUTION: MEET TOASTIE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "The Solution — Meet Toastie",
             "Your AI-powered personal stylist, built into Burnt Toast")

# Big centre statement
add_rect(sl, 0.4, 1.3, 12.5, 1.6, ACCENT)
add_text(sl, "Toastie is not a chatbot.\nToastie is your personal AI stylist — available 24/7, for every customer.",
         0.6, 1.45, 12, 1.2, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cols3 = [
    ("💬", "You Talk.\nToastie Listens.",
     "Tell Toastie your occasion, mood, or style — just like texting a friend."),
    ("✨", "Toastie Builds\nYour Complete Look.",
     "Instant, curated outfit — top, bottom, shoes, bag, accessories. All matched."),
    ("🛍️", "You Shop\nThe Whole Look.",
     "Every item is clickable, swappable, and ready to add to cart — instantly."),
]

for i, (icon, title, body) in enumerate(cols3):
    x = 0.5 + i * 4.3
    card(sl, x, 3.2, 4.0, 3.8, WHITE)
    add_rect(sl, x, 3.2, 4.0, 0.08, ACCENT)
    add_text(sl, icon, x, 3.35, 4.0, 0.8, font_size=34, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.1, 4.2, 3.8, 0.9, font_size=15, bold=True,
             color=INK, align=PP_ALIGN.CENTER)
    add_text(sl, body, x+0.1, 5.2, 3.8, 1.5, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # step arrow
    if i < 2:
        add_text(sl, "→", x+3.95, 4.5, 0.5, 0.5, font_size=24, bold=True,
                 color=ACCENT, align=PP_ALIGN.CENTER)

add_text(sl, "Zero style knowledge required. Zero waiting time. Complete outfit in under 10 seconds.",
         0.5, 7.05, 12.3, 0.35, font_size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (SIMPLE FLOW)
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "How It Works — In 3 Simple Steps",
             "From first message to complete look — faster than browsing a single category page")

# Flow diagram
steps = [
    ("STEP 1", "Tell Toastie", "💬",
     "Type your occasion, mood, or style\n\"Brunch look for this weekend\"\n\"Office outfit — smart casual\"\n\"Freshers night — something fun\"",
     ACCENT),
    ("STEP 2", "Get Your Look", "✨",
     "Toastie instantly builds a\ncomplete styled outfit —\ntop, bottom, shoes, bag +\naccessories, all matched.",
     SAGE),
    ("STEP 3", "Swap, Save & Shop", "🛍️",
     "Don't love a piece? Swap it.\nWant different shoes? Just ask.\nLove the look? Add all to cart\nwith one tap.",
     GOLD),
]

for i, (step_label, title, icon, body, color) in enumerate(steps):
    x = 0.4 + i * 4.3
    # card
    card(sl, x, 1.3, 4.1, 5.6, WHITE)
    # top color strip
    add_rect(sl, x, 1.3, 4.1, 0.9, color)
    add_text(sl, step_label, x+0.1, 1.35, 2, 0.5, font_size=11, bold=True, color=WHITE)
    add_text(sl, icon, x+2.5, 1.35, 1.5, 0.7, font_size=28, align=PP_ALIGN.CENTER)
    # title
    add_text(sl, title, x+0.1, 2.35, 3.9, 0.6, font_size=20, bold=True, color=INK)
    add_divider(sl, x+0.2, 3.05, 3.5, color, 0.03)
    # body
    add_text(sl, body, x+0.2, 3.2, 3.7, 2.5, font_size=13, color=MUTED)
    # step number bottom
    add_text(sl, f"0{i+1}", x+3.3, 6.5, 0.7, 0.5, font_size=32, bold=True,
             color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

    if i < 2:
        add_text(sl, "→", x+4.05, 3.6, 0.4, 0.6, font_size=28, bold=True,
                 color=INK, align=PP_ALIGN.CENTER)

add_rect(sl, 0, 7.1, 13.33, 0.4, INK)
add_text(sl, "Also works with IMAGE UPLOAD — snap a product, Toastie builds the whole outfit around it.",
         0.5, 7.15, 12.3, 0.3, font_size=11, color=CREAM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)
add_rect(sl, 0, 0, 13.33, 1.1, ACCENT)
add_text(sl, "Key Features", 0.5, 0.18, 8, 0.55, font_size=26, bold=True, color=WHITE)
add_text(sl, "Everything Toastie can do — out of the box, today",
         0.5, 0.72, 9, 0.35, font_size=13, color=WHITE, italic=True)
add_text(sl, "BURNT TOAST  ×  TATA", 10.2, 0.3, 3, 0.35,
         font_size=9, color=RGBColor(0xCC,0x88,0x60), align=PP_ALIGN.RIGHT, bold=True)

features = [
    ("💬", "AI Chat Styling",
     "Conversational stylist. User types naturally — Toastie understands occasion, mood & style preference."),
    ("📸", "Image-Based Styling",
     "Upload any product image — Toastie builds a complete outfit around what you already own or love."),
    ("🎯", "Occasion Intelligence",
     "Brunch, office, date night, freshers, IPL, airport — Toastie knows exactly what works for every moment."),
    ("🔄", "Real-Time Swaps",
     "Swap any single item — shoes, top, bag — without regenerating the whole outfit. One tap, instant."),
    ("🛍️", "Shop the Look",
     "Every outfit item is a real, in-stock Burnt Toast product with price. Add to cart in one step."),
    ("🎨", "Vibe Customisation",
     "Smart casual, Y2K, streetwear, minimal — Toastie matches the aesthetic, not just the occasion."),
]

for i, (icon, title, body) in enumerate(features):
    row = i // 3
    col = i % 3
    x = 0.4 + col * 4.3
    y = 1.35 + row * 2.85
    card(sl, x, y, 4.1, 2.6, RGBColor(0x26, 0x26, 0x26))
    add_rect(sl, x, y, 0.6, 2.6, ACCENT)
    add_text(sl, icon, x, y+0.8, 0.6, 0.8, font_size=22, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(sl, title, x+0.7, y+0.15, 3.2, 0.5, font_size=14, bold=True,
             color=ACCENT_LITE)
    add_text(sl, body, x+0.7, y+0.75, 3.2, 1.6, font_size=11.5,
             color=RGBColor(0xCC,0xCC,0xCC))

add_text(sl, "All features work on mobile and desktop. No app download required.",
         0.4, 7.1, 12.5, 0.3, font_size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — GEN Z: THE TARGET AUDIENCE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "The Gen Z Opportunity",
             "Why this audience — and why now")

# Left: stats panel
add_rect(sl, 0.4, 1.3, 5.8, 5.8, ACCENT)
stats = [
    ("₹5.4 Lakh Cr", "Gen Z fashion market in India by 2030"),
    ("73%", "of Gen Z prefer brands that personalise the experience"),
    ("3×", "more likely to buy a full outfit vs. single product"),
    ("60%", "start their shopping journey on mobile"),
    ("#1", "reason for purchase: style inspiration from friends / influencers"),
]
add_text(sl, "Gen Z\nBy The\nNumbers", 0.6, 1.45, 5.4, 1.8,
         font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_divider(sl, 0.7, 3.3, 4.8, WHITE, 0.03)
for i, (stat, label) in enumerate(stats):
    y = 3.5 + i * 0.88
    add_text(sl, stat, 0.7, y, 2.2, 0.6, font_size=19, bold=True, color=WHITE)
    add_text(sl, label, 2.8, y+0.08, 3.2, 0.5, font_size=10.5, color=RGBColor(0xFF,0xDD,0xCC))

# Right: insight cards
insights = [
    ("They shop by vibe, not by category.",
     "Gen Z doesn't think \"I need a top\". They think \"I need a brunch look.\" Toastie speaks their language."),
    ("They want a conversation, not a catalogue.",
     "They message friends for style advice. Toastie is that friend — available inside Burnt Toast, 24/7."),
    ("They share looks, not products.",
     "A full outfit is shareable content. A single T-shirt isn't. Toastie creates shareable moments."),
]
for i, (title, body) in enumerate(insights):
    y = 1.3 + i * 2.05
    card(sl, 6.5, y, 6.5, 1.85, WHITE)
    add_rect(sl, 6.5, y, 0.08, 1.85, ACCENT)
    add_text(sl, title, 6.7, y+0.15, 6.1, 0.5, font_size=14, bold=True, color=INK)
    add_text(sl, body, 6.7, y+0.65, 6.1, 1.0, font_size=12, color=MUTED)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — BUSINESS IMPACT: ENGAGEMENT
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "Business Impact — Customer Engagement",
             "Toastie keeps customers on-site longer and exploring more")

add_rect(sl, 0.4, 1.3, 12.5, 1.1, INK)
add_text(sl, "More time on site  ·  More products discovered  ·  More reasons to come back",
         0.5, 1.55, 12.2, 0.6, font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

metrics = [
    ("4×", "More Products\nViewed Per Session",
     "An outfit recommendation exposes customers to 5–7 products at once vs. browsing 1 at a time.", ACCENT),
    ("↑ 65%", "Longer Time\nOn Site",
     "Conversational shopping is engaging. Customers who interact with the stylist stay significantly longer.", SAGE),
    ("↑ 40%", "Repeat Visit\nRate",
     "\"What should I wear to X?\" brings users back. Toastie gives them a reason to return every time.", GOLD),
    ("↑ 3×", "Category\nDiscovery",
     "Users who came for a dress discover shoes, bags, and accessories they wouldn't have found otherwise.", ACCENT_LITE),
]

for i, (stat, title, body, color) in enumerate(metrics):
    x = 0.4 + i * 3.25
    card(sl, x, 2.7, 3.1, 4.4, WHITE)
    add_rect(sl, x, 2.7, 3.1, 0.06, color)
    add_text(sl, stat, x+0.1, 2.85, 2.9, 0.9, font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.1, 3.85, 2.9, 0.65, font_size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_divider(sl, x+0.4, 4.6, 2.3, LIGHT_GREY, 0.025)
    add_text(sl, body, x+0.15, 4.75, 2.8, 1.9, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_text(sl, "Higher engagement = lower bounce rate = more opportunities to convert.",
         0.4, 7.1, 12.5, 0.35, font_size=12, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS IMPACT: CONVERSION & REVENUE
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)
add_rect(sl, 0, 0, 13.33, 1.1, SAGE)
add_text(sl, "Business Impact — Conversion & Revenue", 0.5, 0.18, 10, 0.55,
         font_size=26, bold=True, color=WHITE)
add_text(sl, "How Toastie directly grows the top line",
         0.5, 0.72, 9, 0.35, font_size=13, color=WHITE, italic=True)
add_text(sl, "BURNT TOAST  ×  TATA", 10.2, 0.3, 3, 0.35,
         font_size=9, color=MUTED, align=PP_ALIGN.RIGHT, bold=True)

# Left column: the numbers
levers = [
    ("Higher Basket Size",
     "A complete outfit = 5–7 items per transaction.\nVs. average of 1.4 items without styling guidance."),
    ("Lower Return Rate",
     "When customers buy a matched outfit, they're\nmore confident — fewer returns, lower costs."),
    ("Upsell Built In",
     "Toastie naturally recommends accessories and\nadd-ons — the stylist equivalent of \"would you\nlike fries with that?\""),
    ("New Customer Acquisition",
     "\"The brand that styled me perfectly\" drives\nword-of-mouth and social sharing among Gen Z."),
]

for i, (title, body) in enumerate(levers):
    y = 1.4 + i * 1.45
    add_rect(sl, 0.4, y, 0.06, 1.2, SAGE)
    add_text(sl, title, 0.65, y+0.05, 5.5, 0.45, font_size=14, bold=True, color=RGBColor(0xAA,0xCC,0xAA))
    add_text(sl, body, 0.65, y+0.5, 5.5, 0.85, font_size=11.5, color=RGBColor(0xCC,0xCC,0xCC))

# Right column: big impact numbers
add_rect(sl, 7.2, 1.3, 5.8, 5.8, RGBColor(0x22,0x22,0x22))
big_stats = [
    ("2.5×", "Average Order\nValue Uplift", SAGE),
    ("↓ 25%", "Reduction in\nCart Abandonment", GOLD),
    ("↑ 35%", "Conversion Rate\nImprovement", ACCENT_LITE),
]
add_text(sl, "Projected Impact", 7.5, 1.5, 5.2, 0.5, font_size=14, bold=True,
         color=MUTED, align=PP_ALIGN.CENTER)
add_divider(sl, 7.7, 2.05, 4.8, SAGE, 0.03)
for i, (num, label, color) in enumerate(big_stats):
    y = 2.3 + i * 1.7
    add_rect(sl, 7.4, y, 5.2, 1.45, RGBColor(0x2C,0x2C,0x2C))
    add_text(sl, num, 7.5, y+0.1, 3.5, 0.9, font_size=38, bold=True, color=color)
    add_text(sl, label, 10.2, y+0.2, 2.3, 0.9, font_size=12, color=RGBColor(0xCC,0xCC,0xCC))

add_text(sl, "Based on industry benchmarks for AI-assisted styling tools in fashion e-commerce.",
         0.4, 7.15, 12.5, 0.3, font_size=9, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — TOASTIE VS REGULAR CHATBOT
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "Toastie vs. A Regular Chatbot",
             "This is a fundamentally different kind of AI — not just a FAQ bot")

headers = ["Feature", "Regular Chatbot", "Toastie — AI Stylist"]
rows = [
    ("Purpose",          "Answer FAQ / support tickets",       "Build personalised outfit recommendations"),
    ("Understands",      "Keywords & simple commands",         "Occasion, vibe, body type, style preference"),
    ("Output",           "Text reply",                         "Complete styled outfit with real products"),
    ("Products shown",   "1 item at a time (if at all)",       "5–7 matched items per outfit, instantly"),
    ("Personalisation",  "None / rule-based",                  "AI-powered, learns from the conversation"),
    ("Image input",      "❌ Not supported",                   "✅ Upload any product image"),
    ("Business value",   "Low — reduces support load only",    "High — drives sales, AOV, & discovery"),
    ("Gen Z appeal",     "❌ Feels robotic and generic",       "✅ Feels like texting a stylish friend"),
]

col_widths = [2.8, 4.1, 5.5]
col_starts = [0.4, 3.3, 7.5]
row_h = 0.62

# Header row
for j, (hdr, w, x) in enumerate(zip(headers, col_widths, col_starts)):
    bg = INK if j == 0 else (MUTED if j == 1 else ACCENT)
    add_rect(sl, x, 1.3, w, 0.55, bg)
    add_text(sl, hdr, x+0.1, 1.35, w-0.2, 0.45, font_size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

for i, (feature, chatbot_val, toastie_val) in enumerate(rows):
    y = 1.85 + i * row_h
    bg_row = WHITE if i % 2 == 0 else CREAM_SOFT
    for j, (val, w, x) in enumerate(zip([feature, chatbot_val, toastie_val], col_widths, col_starts)):
        add_rect(sl, x, y, w, row_h, bg_row)
        c = INK if j == 0 else (MUTED if j == 1 else SAGE)
        bold = j == 0
        add_text(sl, val, x+0.1, y+0.1, w-0.2, row_h-0.15, font_size=11.5, bold=bold, color=c)

# line dividers
for x in col_starts[1:]:
    add_rect(sl, x-0.02, 1.3, 0.02, len(rows)*row_h+0.55, LIGHT_GREY)

add_text(sl, "Toastie is a revenue driver. A chatbot is a cost centre.",
         0.4, 7.1, 12.5, 0.35, font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — FUTURE ROADMAP
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)
add_rect(sl, 0, 0, 13.33, 1.1, GOLD)
add_text(sl, "Future Possibilities & Scalability", 0.5, 0.18, 10, 0.55,
         font_size=26, bold=True, color=INK)
add_text(sl, "Today's MVP is just the beginning — here's where Toastie can go",
         0.5, 0.72, 9, 0.35, font_size=13, color=INK, italic=True)
add_text(sl, "BURNT TOAST  ×  TATA", 10.2, 0.3, 3, 0.35,
         font_size=9, color=INK, align=PP_ALIGN.RIGHT, bold=True)

phases = [
    ("Phase 1\n(Now)", "Core Stylist", ACCENT, [
        "AI chat + image styling",
        "Complete outfit builder",
        "Swap & customise any item",
        "Occasion intelligence",
        "Real-time catalogue",
    ]),
    ("Phase 2\n(3–6 months)", "Personalisation", SAGE, [
        "User style profiles & history",
        "\"My Wardrobe\" feature",
        "Wishlist + saved looks",
        "Toastie learns your taste",
        "Push notification styling",
    ]),
    ("Phase 3\n(6–12 months)", "Social & Commerce", GOLD, [
        "Share looks to Instagram/WhatsApp",
        "Influencer style collections",
        "\"Get Celebrity's Look\" feature",
        "Group styling (friends shop together)",
        "Loyalty points per look shopped",
    ]),
    ("Phase 4\n(12+ months)", "Next Gen Tech", RGBColor(0x7B,0x68,0xEE), [
        "Virtual try-on (AR)",
        "Size & fit recommendation AI",
        "Video lookbook generation",
        "Cross-brand TATA styling",
        "B2B licensing to other brands",
    ]),
]

for i, (phase, title, color, items) in enumerate(phases):
    x = 0.4 + i * 3.25
    card(sl, x, 1.35, 3.1, 5.8, RGBColor(0x22,0x22,0x22))
    add_rect(sl, x, 1.35, 3.1, 1.1, color)
    add_text(sl, phase, x+0.1, 1.4, 2.9, 0.55, font_size=11, bold=True,
             color=INK if color==GOLD else WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.1, 1.9, 2.9, 0.45, font_size=14, bold=True,
             color=INK if color==GOLD else WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = 2.65 + j * 0.82
        add_rect(sl, x+0.25, y+0.18, 0.12, 0.12, color)
        add_text(sl, item, x+0.5, y+0.05, 2.5, 0.65, font_size=11,
                 color=RGBColor(0xCC,0xCC,0xCC))

add_text(sl, "Each phase is modular — add features as the business grows, without rebuilding from scratch.",
         0.4, 7.15, 12.5, 0.3, font_size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "The Market Opportunity",
             "Why this is the right product, at the right time, for the right brand")

# Left: big numbers
add_rect(sl, 0.4, 1.3, 5.8, 5.8, ACCENT)
market_stats = [
    ("₹1.8 Lakh Cr", "Indian fashion\ne-commerce market (2025)"),
    ("28%", "Year-on-year growth\nin online fashion"),
    ("50 Cr+", "Gen Z consumers\nin India today"),
    ("45%", "Of Gen Z already\nshops fashion online"),
]
add_text(sl, "The\nOpportunity", 0.65, 1.5, 5.2, 1.3, font_size=26, bold=True, color=WHITE)
add_divider(sl, 0.7, 2.9, 4.8, WHITE, 0.03)
for i, (num, label) in enumerate(market_stats):
    y = 3.1 + i * 0.97
    card(sl, 0.65, y, 5.2, 0.78, RGBColor(0xCC,0x55,0x30))
    add_text(sl, num, 0.8, y+0.1, 2.3, 0.6, font_size=20, bold=True, color=WHITE)
    add_text(sl, label, 3.1, y+0.12, 2.5, 0.55, font_size=10.5, color=RGBColor(0xFF,0xEE,0xDD))

# Right: why TATA / Burnt Toast wins
add_text(sl, "Why Burnt Toast\nis Perfectly Positioned", 6.7, 1.35, 6.3, 0.8,
         font_size=18, bold=True, color=INK)
add_divider(sl, 6.7, 2.25, 6.1, ACCENT, 0.04)

advantages = [
    ("🏆 First Mover",
     "No Indian fast-fashion brand has an AI personal stylist. Burnt Toast gets there first."),
    ("🎯 Gen Z Native",
     "The brand, the aesthetic, the price point — everything about Burnt Toast is built for Gen Z."),
    ("🔗 TATA Ecosystem",
     "Access to TATA's data infrastructure, brand trust, and cross-brand synergies (Zudio, Westside, etc.)"),
    ("📈 Proven Model",
     "Global brands like Stitch Fix & ASOS AI have shown that styling AI dramatically lifts revenue."),
]
for i, (title, body) in enumerate(advantages):
    y = 2.5 + i * 1.2
    add_text(sl, title, 6.7, y, 6.3, 0.4, font_size=13, bold=True, color=ACCENT)
    add_text(sl, body, 6.7, y+0.4, 6.3, 0.65, font_size=11.5, color=MUTED)
    if i < 3:
        add_divider(sl, 6.7, y+1.1, 6.1, LIGHT_GREY, 0.02)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — THE CUSTOMER EXPERIENCE JOURNEY
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)
add_rect(sl, 0, 0, 13.33, 1.1, SAGE)
add_text(sl, "The Customer Experience Journey", 0.5, 0.18, 10, 0.55,
         font_size=26, bold=True, color=WHITE)
add_text(sl, "From discovery to purchase — how Toastie transforms every step",
         0.5, 0.72, 9, 0.35, font_size=13, color=WHITE, italic=True)
add_text(sl, "BURNT TOAST  ×  TATA", 10.2, 0.3, 3, 0.35,
         font_size=9, color=INK, align=PP_ALIGN.RIGHT, bold=True)

journey = [
    ("🔍", "Discovers BT", "Via Instagram,\nfriend's reco,\nor TATA app"),
    ("💬", "Chats with Toastie", "Types occasion\nor uploads a\nproduct photo"),
    ("✨", "Gets Styled", "Complete outfit\nin 5 seconds\nwith real products"),
    ("🔄", "Customises", "Swaps items,\nchanges vibe,\nasked follow-ups"),
    ("🛒", "Adds to Cart", "One tap to add\nthe full look\nto cart"),
    ("📤", "Shares the Look", "Screenshots &\nshares with friends\n→ new customers"),
]

# Timeline line
add_rect(sl, 0.5, 3.8, 12.3, 0.06, SAGE)
for i, (icon, title, sub) in enumerate(journey):
    x = 0.55 + i * 2.17
    # dot on timeline
    add_rect(sl, x+0.4, 3.65, 0.32, 0.32, ACCENT)
    # card above/below alternating
    if i % 2 == 0:
        y_card = 1.4
        y_arrow = 3.55
    else:
        y_card = 4.3
        y_arrow = 3.9

    card(sl, x, y_card, 2.0, 1.9, RGBColor(0x26,0x26,0x26))
    add_text(sl, icon, x, y_card+0.1, 2.0, 0.6, font_size=24, align=PP_ALIGN.CENTER)
    add_text(sl, title, x+0.05, y_card+0.7, 1.9, 0.45, font_size=11, bold=True,
             color=SAGE, align=PP_ALIGN.CENTER)
    add_text(sl, sub, x+0.05, y_card+1.15, 1.9, 0.7, font_size=10,
             color=RGBColor(0xAA,0xAA,0xAA), align=PP_ALIGN.CENTER)

add_rect(sl, 0, 7.1, 13.33, 0.4, RGBColor(0x20,0x20,0x20))
add_text(sl, "Every step of the journey is an opportunity for Burnt Toast to delight, retain and grow.",
         0.4, 7.15, 12.5, 0.3, font_size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — DEVELOPMENT ROADMAP
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, CREAM_SOFT)
slide_header(sl, "Development Roadmap — From Demo to Launch",
             "Where we are today and what happens next with your approval")

# Status indicator
add_rect(sl, 0.4, 1.35, 12.5, 0.75, GOLD)
add_text(sl, "📍  CURRENT STATUS: Testing & Demo Phase — Fully Functional Product Ready for Review",
         0.6, 1.5, 12.2, 0.45, font_size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)

phases_road = [
    ("✅ DONE", "Phase 0\nFoundation", SAGE,
     ["AI stylist (Toastie) live", "Catalogue: 280+ products", "Chat + Image styling",
      "Mobile-responsive UI", "Rate limiting & security"]),
    ("🟡 NOW", "Phase 1\nDemo & Approval", GOLD,
     ["Panel presentation", "Live demo walkthrough", "Feedback collection",
      "Business case approval", "Go / No-Go decision"]),
    ("🔵 NEXT", "Phase 2\nProduction Ready", ACCENT,
     ["Performance optimisation", "Full QA & testing", "Analytics dashboard",
      "User accounts & profiles", "App Store listing"]),
    ("⭕ FUTURE", "Phase 3\nGrowth Features", MUTED,
     ["Virtual try-on (AR)", "Social sharing tools", "Loyalty integration",
      "Influencer collections", "Cross-TATA expansion"]),
]

for i, (status, phase, color, items) in enumerate(phases_road):
    x = 0.4 + i * 3.25
    card(sl, x, 2.35, 3.1, 4.85, WHITE)
    add_rect(sl, x, 2.35, 3.1, 0.75, color)
    add_text(sl, status, x+0.1, 2.38, 2.9, 0.4, font_size=11, bold=True,
             color=WHITE if color != GOLD else INK, align=PP_ALIGN.CENTER)
    add_text(sl, phase, x+0.1, 2.72, 2.9, 0.35, font_size=12, bold=True,
             color=WHITE if color != GOLD else INK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = 3.25 + j * 0.72
        prefix = "✓ " if status == "✅ DONE" else "→ "
        add_text(sl, prefix + item, x+0.2, y, 2.8, 0.6, font_size=11.5, color=MUTED)

add_rect(sl, 0, 7.1, 13.33, 0.4, INK)
add_text(sl, "With your approval today → Phase 2 begins within 2 weeks.",
         0.4, 7.15, 12.5, 0.3, font_size=12, bold=True, color=ACCENT_LITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING / THE ASK
# ═══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, INK)
add_rect(sl, 0, 0, 5.0, 7.5, ACCENT)

# Left
add_text(sl, "The\nVision", 0.4, 0.8, 4.2, 1.8, font_size=40, bold=True, color=WHITE)
add_divider(sl, 0.4, 2.75, 4.0, WHITE, 0.04)
add_text(sl, "Every Burnt Toast\ncustomer walks away\nfeeling styled,\nconfident, and\ninspired —\nevery single time.",
         0.4, 2.95, 4.3, 3.2, font_size=18, color=WHITE)
add_text(sl, "Powered by AI.\nBuilt for Gen Z.\nDriven by TATA.", 0.4, 6.2, 4.3, 1.1,
         font_size=13, color=RGBColor(0xFF,0xDD,0xCC), italic=True)

# Right
add_text(sl, "What We're Asking For", 5.4, 0.7, 7.5, 0.6,
         font_size=22, bold=True, color=WHITE)
add_divider(sl, 5.4, 1.4, 7.3, ACCENT, 0.04)

asks = [
    ("✅", "Approval to move from Demo → Production"),
    ("📋", "Alignment on Phase 2 feature priorities"),
    ("💰", "Budget sign-off for full development"),
    ("🤝", "TATA ecosystem integration discussions"),
    ("📅", "Target launch date agreement"),
]
for i, (icon, ask) in enumerate(asks):
    y = 1.65 + i * 0.9
    add_rect(sl, 5.4, y, 7.3, 0.75, RGBColor(0x28,0x28,0x28))
    add_text(sl, icon, 5.5, y+0.12, 0.6, 0.5, font_size=18, align=PP_ALIGN.CENTER)
    add_text(sl, ask, 6.15, y+0.15, 6.4, 0.45, font_size=13.5, color=WHITE)

# Final tagline
add_rect(sl, 5.2, 6.2, 7.7, 1.0, RGBColor(0x22,0x22,0x22))
add_text(sl, "\"The brand that styles you is the brand you stay loyal to.\"",
         5.4, 6.35, 7.3, 0.65, font_size=14, italic=True, color=ACCENT_LITE, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(sl, 0, 7.1, 13.33, 0.4, RGBColor(0x20,0x20,0x20))
add_text(sl, "BURNT TOAST  ×  TATA  |  AI Fashion Stylist  |  Confidential  |  2026",
         0.4, 7.15, 12.5, 0.3, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────
out = "Burnt_Toast_Pitch_Deck.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
